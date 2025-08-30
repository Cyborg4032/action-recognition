
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np
import io

# Create presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sign Language Detection"
subtitle.text = "Using MediaPipe and LSTM"

# Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Project Overview"
content.text = ("- Real-time sign language detection system\n"
                "- Recognizes gestures: hello, thanks, I love you\n" 
                "- Uses MediaPipe for feature extraction\n"
                "- LSTM neural network for classification")

# MediaPipe Holistic
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "MediaPipe Holistic"
content.text = ("- Detects 33 pose, 21 hand, and 468 face landmarks\n"
                "- Provides 3D coordinates of each landmark\n"
                "- Processes video frames in real-time")

# Data Collection
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Data Collection"
content.text = ("- Collected 30 sequences for each gesture\n"
                "- Each sequence contains 30 frames\n"
                "- Saved as numpy arrays of keypoints\n"
                "- Total of 1662 features per frame")

# Model Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "LSTM Model Architecture"
content.text = ("- Input shape: (30, 1662)\n"
                "- 3 LSTM layers (64, 128, 64 units)\n"
                "- 3 Dense layers (64, 32, 3 units)\n"
                "- Softmax activation for classification\n"
                "- Adam optimizer with 2000 epochs")

# Training Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Training Results"
content.text = ("- Achieved 95%+ accuracy on test set\n"
                "- Confusion matrix shows good separation\n"
                "- Model saved as action.h5 for deployment")

# Real-Time Detection
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Real-Time Detection"
content.text = ("- Processes webcam feed at 30fps\n"
                "- Displays detected gesture and confidence\n"
                "- Visualizes keypoints on screen\n"
                "- Threshold of 0.8 for prediction")

# Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion & Future Work"
content.text = ("- Successful proof of concept\n"
                "- Potential applications in accessibility\n"
                "- Future work:\n"
                "  - Expand vocabulary\n"
                "  - Improve robustness\n"
                "  - Mobile deployment")

# Save presentation
prs.save('Sign_Language_Detection_Presentation.pptx')
